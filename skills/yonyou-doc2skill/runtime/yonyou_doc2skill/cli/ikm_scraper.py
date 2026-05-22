#!/usr/bin/env python3
"""IKM knowledge map to skill converter."""

import argparse
import json
import os
import re
import zipfile
from xml.etree import ElementTree
from pathlib import Path
from typing import Any
from urllib.parse import parse_qs, unquote_plus, urlparse

import requests
from bs4 import BeautifulSoup

from yonyou_doc2skill.cli.arguments.ikm import add_ikm_arguments
from yonyou_doc2skill.cli.skill_converter import SkillConverter


def _safe_filename(value: str, fallback: str = "item") -> str:
    safe = re.sub(r"[^\w\-.一-龥]+", "_", value.strip(), flags=re.UNICODE)
    safe = safe.strip("._")
    return safe or fallback


def _strip_html(value: Any) -> str:
    text = str(value or "")
    if not text:
        return ""
    if "<" not in text or ">" not in text:
        return text.strip()
    return BeautifulSoup(text, "html.parser").get_text("", strip=True)


def _detail_extension(detail: dict[str, Any]) -> str:
    detail_type = str(detail.get("detailtype") or "").lower().strip()
    if detail_type:
        return detail_type.lstrip(".")
    code = str(detail.get("detailTypeCode") or "").upper()
    return {
        "TC": "pdf",
        "TA": "pptx",
        "TB": "docx",
        "TD": "xlsx",
        "TXT": "txt",
    }.get(code, "bin")


PARSEABLE_ATTACHMENT_EXTENSIONS = {
    "txt",
    "md",
    "csv",
    "json",
    "xml",
    "log",
    "html",
    "htm",
    "pdf",
    "docx",
    "pptx",
}


class IKMToSkillConverter(SkillConverter):
    """Extract IKM knowledge maps into a local skill."""

    SOURCE_TYPE = "ikm"

    def __init__(self, config: dict[str, Any]):
        super().__init__(config)
        self.base_url = str(config.get("base_url") or "https://ikm.yonyou.com").rstrip("/")
        self.mode = config.get("mode") or "map"
        self.url = config.get("url")
        self.cookie = config.get("cookie") or os.getenv("IKM_COOKIE")
        self.pk = config.get("pk")
        self.keyword = config.get("keyword")
        self.actionlocid = config.get("actionlocid")
        self.actionloctype = config.get("actionloctype") or "portal"
        self.asset_module = config.get("module")
        self.asset_moduleobjid = config.get("moduleobjid")
        self.asset_moduleobjname = config.get("moduleobjname")
        if self.mode == "asset" and self.url:
            self._apply_asset_url(self.url)
        self.portal_endpoint = config.get("portal_endpoint") or "/space/initSharePortalChannel"
        self.search_endpoint = config.get("search_endpoint") or "/asset/getAssetsByEs"
        self.max_assets = int(config.get("max_assets") or 100)
        self.parse_attachments = bool(config.get("parse_attachments"))
        self.download_all_attachments = bool(config.get("download_attachments"))
        self.download_attachments = self.download_all_attachments or self.parse_attachments
        self.max_attachment_chars = int(config.get("max_attachment_chars") or 12000)
        self.data_file = f"output/{self.name}_extracted.json"
        self.session = config.get("_session") or requests.Session()
        self.extracted_data: dict[str, Any] | None = None
        self._current_map_name = config.get("moduleobjname") or config.get("name") or ""

        if not self.cookie and not config.get("from_json"):
            raise ValueError("IKM cookie is required. Use --cookie or set IKM_COOKIE.")
        if self.mode not in {"map", "portal", "search", "asset"}:
            raise ValueError("IKM mode must be one of: map, portal, search, asset.")
        if self.mode == "map" and not self.pk and not config.get("from_json"):
            raise ValueError("IKM map pk is required. Use --pk.")
        if self.mode == "search" and not self.keyword and not config.get("from_json"):
            raise ValueError("IKM search keyword is required. Use --keyword.")
        if self.mode == "asset" and not self.pk and not config.get("from_json"):
            raise ValueError("IKM asset pk is required. Use --pk or --url with pkasset.")
        if self.mode in {"map", "portal", "search", "asset"} and not self.actionlocid and not config.get("from_json"):
            raise ValueError("IKM actionlocid is required. Use --actionlocid.")

        self.session.headers.update(
            {
                "Cookie": self.cookie,
                "User-Agent": "Mozilla/5.0",
                "X-Requested-With": "XMLHttpRequest",
            }
        )

    def _apply_asset_url(self, url: str) -> None:
        parsed = urlparse(url)
        query = self._parse_raw_query(parsed.query)
        self.pk = self.pk or query.get("pkasset") or query.get("pkKnowledgeAsset")
        self.actionlocid = self.actionlocid or query.get("actionlocid")
        self.actionloctype = query.get("actionloctype") or self.actionloctype

        path_payload = query.get("path")
        if path_payload:
            decoded_path = self._decode_url_path_payload(path_payload)
            if decoded_path:
                self.asset_module = self.asset_module or decoded_path.get("module")
                self.asset_moduleobjid = self.asset_moduleobjid or decoded_path.get("moduleobjid")
                self.asset_moduleobjname = (
                    self.asset_moduleobjname or decoded_path.get("moduleobjname")
                )

        if not self.asset_module:
            fragment = parsed.fragment.lower()
            self.asset_module = "search" if "search" in fragment else "asset"
        if self.asset_moduleobjid is None:
            self.asset_moduleobjid = ""
        if self.asset_moduleobjname is None:
            self.asset_moduleobjname = ""

    def _parse_raw_query(self, query: str) -> dict[str, str]:
        parsed: dict[str, str] = {}
        for part in query.split("&"):
            if not part:
                continue
            key, _, raw_value = part.partition("=")
            parsed[unquote_plus(key)] = unquote_plus(raw_value)
        return parsed

    def _decode_url_path_payload(self, payload: str) -> dict[str, Any]:
        import base64

        normalized = payload.replace(" ", "+")
        padding = "=" * ((4 - len(normalized) % 4) % 4)
        try:
            decoded = base64.b64decode(normalized + padding).decode("utf-8")
            return json.loads(decoded)
        except (ValueError, json.JSONDecodeError):
            return {}

    def _post_json(self, path: str, data: dict[str, Any]) -> dict[str, Any]:
        response = self.session.post(f"{self.base_url}{path}", data=data, timeout=30)
        response.raise_for_status()
        payload = response.json()
        if str(payload.get("code")) != "200":
            raise RuntimeError(f"IKM API failed for {path}: {payload.get('msg') or payload}")
        return payload

    def _get_json(self, path: str, params: dict[str, Any]) -> dict[str, Any]:
        response = self.session.get(f"{self.base_url}{path}", params=params, timeout=30)
        response.raise_for_status()
        payload = response.json()
        if str(payload.get("code")) != "200":
            raise RuntimeError(f"IKM API failed for {path}: {payload.get('msg') or payload}")
        return payload

    def extract(self):
        if self.config.get("from_json"):
            self.load_extracted_data(self.config["from_json"])
            return

        if self.mode == "map":
            result = self._extract_map()
        elif self.mode == "portal":
            result = self._extract_portal()
        elif self.mode == "search":
            result = self._extract_search()
        elif self.mode == "asset":
            result = self._extract_asset()
        else:
            raise ValueError(f"Unsupported IKM mode: {self.mode}")

        self.extracted_data = result

        Path("output").mkdir(exist_ok=True)
        Path(self.data_file).write_text(
            json.dumps(result, ensure_ascii=False, indent=2), encoding="utf-8"
        )

    def _extract_map(self) -> dict[str, Any]:
        map_payload = self._post_json(
            "/knowledgemap/getKnowledgeMapByPK",
            {
                "pkMapview": self.pk,
                "actionloctype": self.actionloctype,
                "actionlocid": self.actionlocid,
            },
        )
        map_info = map_payload["data"]
        self._current_map_name = map_info.get("mapname") or self._current_map_name

        assets_payload = self._post_json(
            "/knowledgemap/getMapData",
            {
                "pageIndex": 1,
                "pageSize": self.max_assets,
                "catecode": "",
                "objecttype": "asset",
                "order": "sortorder",
                "pkMapView": self.pk,
                "actionloctype": self.actionloctype,
                "actionlocid": self.actionlocid,
                "keyword": "",
            },
        )
        assets_data = assets_payload.get("data") or {}
        assets = list(assets_data.get("pageList") or [])[: self.max_assets]

        for asset in assets:
            asset["attachments"] = self._fetch_attachments(asset)

        return {
            "source_type": "ikm",
            "mode": self.mode,
            "base_url": self.base_url,
            "actionloctype": self.actionloctype,
            "actionlocid": self.actionlocid,
            "map": map_info,
            "categories": assets_data.get("mapCates") or [],
            "assets": assets,
            "total_assets": len(assets),
        }

    def _extract_portal(self) -> dict[str, Any]:
        payload = self._post_json(
            self.portal_endpoint,
            {
                "tenantid": "yonyou",
                "actionloctype": self.actionloctype,
                "actionlocid": self.actionlocid,
                "actionboardid": "",
                "sortType": "creattime",
            },
        )
        assets = self._assets_from_payload_data(payload)
        for asset in assets:
            asset["attachments"] = self._fetch_attachments(asset)

        return {
            "source_type": "ikm",
            "mode": self.mode,
            "base_url": self.base_url,
            "actionloctype": self.actionloctype,
            "actionlocid": self.actionlocid,
            "portal": {
                "endpoint": self.portal_endpoint,
                "actionlocid": self.actionlocid,
            },
            "assets": assets,
            "total_assets": len(assets),
        }

    def _extract_asset(self) -> dict[str, Any]:
        asset = {
            "pkKnowledgeAsset": self.pk,
            "assetname": self.config.get("assetname") or f"iKM Asset {self.pk}",
            "knowledgeitemdes": self.config.get("description") or "",
            "filenum": 0,
            "source_url": self.url or "",
            "module": self.asset_module or "",
            "moduleobjid": self.asset_moduleobjid or "",
            "moduleobjname": self.asset_moduleobjname or "",
        }
        asset["attachments"] = self._fetch_attachments(asset)
        asset["filenum"] = len(asset["attachments"])
        detail_names = [
            _strip_html(detail.get("oridetailname") or detail.get("detailname"))
            for detail in asset["attachments"]
        ]
        asset["detailnames"] = ", ".join(name for name in detail_names if name)

        return {
            "source_type": "ikm",
            "mode": self.mode,
            "base_url": self.base_url,
            "actionloctype": self.actionloctype,
            "actionlocid": self.actionlocid,
            "asset": asset,
            "assets": [asset],
            "total_assets": 1,
        }

    def _extract_search(self) -> dict[str, Any]:
        payload = self._get_or_post_json(
            self.search_endpoint,
            {
                "assetparam": self.keyword,
                "sorttype": "default",
                "pageIndex": 1,
                "searchdoc": json.dumps(
                    {"type": "", "industry": "", "product": "", "field": "", "filetype": ""},
                    ensure_ascii=False,
                ),
                "searchscope": "asset",
                "exactorpartial": "exact",
                "searchcreattime": "",
                "searchmodifiedtime": "",
                "pkscene": self.actionlocid,
                "behaviorPath": json.dumps(
                    {"module": "search", "moduleobjid": "", "moduleobjname": ""},
                    ensure_ascii=False,
                ),
                "actionloctype": self.actionloctype,
                "actionlocid": self.actionlocid,
                "actionboardid": "",
            },
        )
        data = payload.get("data") or {}
        assets = self._assets_from_payload_data(data)
        for asset in assets:
            asset["attachments"] = self._fetch_attachments(asset)

        return {
            "source_type": "ikm",
            "mode": self.mode,
            "base_url": self.base_url,
            "actionloctype": self.actionloctype,
            "actionlocid": self.actionlocid,
            "keyword": self.keyword,
            "search": {
                "endpoint": self.search_endpoint,
                "keyword": self.keyword,
            },
            "assets": assets,
            "total_assets": len(assets),
        }

    def _assets_from_payload_data(self, data: Any) -> list[dict[str, Any]]:
        if isinstance(data, list):
            return data[: self.max_assets]
        if not isinstance(data, dict):
            return []

        for key in ("pageList", "firstAssetList", "newassetlist", "list", "records", "rows", "data"):
            value = data.get(key)
            if isinstance(value, list):
                return value[: self.max_assets]
            if isinstance(value, dict):
                nested = self._assets_from_payload_data(value)
                if nested:
                    return nested[: self.max_assets]
        return []

    def _get_or_post_json(self, path: str, params: dict[str, Any]) -> dict[str, Any]:
        if path.endswith("getAssetsByEs"):
            return self._get_json(path, params)
        return self._post_json(path, params)

    def _fetch_attachments(self, asset: dict[str, Any]) -> list[dict[str, Any]]:
        asset_pk = asset.get("pkKnowledgeAsset")
        if not asset_pk:
            return []

        try:
            payload = self._get_json(
                "/file/getHlDetail",
                {
                    "assetpk": asset_pk,
                    "param": "",
                    "openauthtype": self.actionloctype,
                    "openlocid": self.actionlocid,
                },
            )
            attachments = list(payload.get("data") or [])
        except Exception as exc:
            asset["attachment_error"] = str(exc)
            return []

        if self.download_attachments:
            for attachment in attachments:
                if self.parse_attachments and not self.download_all_attachments:
                    ext = _detail_extension(attachment).lower()
                    if ext not in PARSEABLE_ATTACHMENT_EXTENSIONS:
                        attachment["download_skipped"] = (
                            f"Skipped non-parseable attachment type: {ext}"
                        )
                        continue
                self._download_attachment(attachment)

        return attachments

    def _download_attachment(self, attachment: dict[str, Any]) -> None:
        detail_pk = attachment.get("pkKnowledgeDetail")
        if not detail_pk:
            return

        ext = _detail_extension(attachment)
        filename = f"{_safe_filename(str(detail_pk))}.{ext}"
        downloads_dir = Path(self.skill_dir) / "downloads"
        downloads_dir.mkdir(parents=True, exist_ok=True)
        path = downloads_dir / filename
        params = self._download_params(detail_pk)

        try:
            response = self.session.get(
                f"{self.base_url}/file/downloadFileSingle", params=params, timeout=120
            )
            response.raise_for_status()
            path.write_bytes(response.content)
            attachment["download_path"] = str(path)
            if self.parse_attachments:
                self._parse_downloaded_attachment(attachment, path)
        except Exception as exc:
            attachment["download_error"] = str(exc)

    def _download_params(self, detail_pk: str) -> dict[str, Any]:
        params = {
            "pkDetatil": detail_pk,
            "actionloctype": self.actionloctype,
            "actionlocid": self.actionlocid,
        }
        if self.mode == "map":
            params.update(
                {
                    "module": "mapview",
                    "moduleobjid": self.pk,
                    "moduleobjname": self._map_name_for_param(),
                }
            )
        elif self.mode == "search":
            params.update({"module": "search", "moduleobjid": "", "moduleobjname": ""})
        elif self.mode == "portal":
            params.update(
                {
                    "module": "portal",
                    "moduleobjid": self.actionlocid,
                    "moduleobjname": self._source_title(),
                }
            )
        elif self.mode == "asset":
            params.update(
                {
                    "module": self.asset_module or "asset",
                    "moduleobjid": self.asset_moduleobjid or "",
                    "moduleobjname": self.asset_moduleobjname or "",
                }
            )
        return params

    def _parse_downloaded_attachment(self, attachment: dict[str, Any], path: Path) -> None:
        """Parse downloaded attachment text into a reference file."""
        detail_pk = str(attachment.get("pkKnowledgeDetail") or path.stem)
        attachment.pop("parse_error", None)
        try:
            text = self._extract_attachment_text(path)
            text = self._normalize_text(text)
            if not text:
                attachment["parse_error"] = "No readable text extracted from attachment."
                return

            if len(text) > self.max_attachment_chars:
                text = text[: self.max_attachment_chars].rstrip() + "\n\n[内容已按字符上限截断]"

            content_dir = Path(self.skill_dir) / "references" / "content"
            content_dir.mkdir(parents=True, exist_ok=True)
            content_path = content_dir / f"{_safe_filename(detail_pk)}.md"
            title = attachment.get("oridetailname") or attachment.get("detailname") or detail_pk
            content_path.write_text(f"# {title}\n\n{text}\n", encoding="utf-8")

            attachment["content_path"] = str(content_path.relative_to(Path(self.skill_dir)))
            attachment["content_chars"] = len(text)
            attachment["content_excerpt"] = text[:500].strip()
            attachment.pop("parse_error", None)
        except Exception as exc:
            attachment["parse_error"] = str(exc)

    def _extract_attachment_text(self, path: Path) -> str:
        ext = path.suffix.lower()
        if ext in {".txt", ".md", ".csv", ".json", ".xml", ".log"}:
            return path.read_text(encoding="utf-8", errors="ignore")
        if ext in {".html", ".htm"}:
            soup = BeautifulSoup(path.read_text(encoding="utf-8", errors="ignore"), "html.parser")
            return soup.get_text("\n", strip=True)
        if ext == ".pdf":
            return self._extract_pdf_text(path)
        if ext == ".docx":
            return self._extract_docx_text(path)
        if ext == ".pptx":
            return self._extract_pptx_text(path)
        raise ValueError(f"Unsupported attachment type for parsing: {ext or 'unknown'}")

    def _extract_pdf_text(self, path: Path) -> str:
        try:
            import fitz
        except ImportError as exc:
            raise RuntimeError("PyMuPDF is required to parse PDF attachments.") from exc

        parts: list[str] = []
        with fitz.open(path) as doc:
            for page_index, page in enumerate(doc, start=1):
                text = page.get_text("text").strip()
                if text:
                    parts.append(f"## Page {page_index}\n\n{text}")
        return "\n\n".join(parts)

    def _extract_docx_text(self, path: Path) -> str:
        try:
            import docx
        except ImportError:
            return self._extract_docx_text_stdlib(path)

        document = docx.Document(path)
        parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
        for table in document.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
        return "\n\n".join(parts)

    def _extract_docx_text_stdlib(self, path: Path) -> str:
        """Extract DOCX text without optional python-docx dependency."""
        namespaces = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
        parts: list[str] = []
        with zipfile.ZipFile(path) as docx_zip:
            names = [
                name
                for name in docx_zip.namelist()
                if name == "word/document.xml"
                or (name.startswith("word/header") and name.endswith(".xml"))
                or (name.startswith("word/footer") and name.endswith(".xml"))
            ]
            for name in names:
                root = ElementTree.fromstring(docx_zip.read(name))
                for paragraph in root.findall(".//w:p", namespaces):
                    texts = [
                        node.text or ""
                        for node in paragraph.findall(".//w:t", namespaces)
                        if node.text
                    ]
                    text = "".join(texts).strip()
                    if text:
                        parts.append(text)
        return "\n\n".join(parts)

    def _extract_pptx_text(self, path: Path) -> str:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise RuntimeError("python-pptx is required to parse PPTX attachments.") from exc

        presentation = Presentation(path)
        parts: list[str] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            slide_parts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_parts.append(shape.text.strip())
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            slide_parts.append(" | ".join(cells))
            if slide_parts:
                parts.append(f"## Slide {slide_index}\n\n" + "\n\n".join(slide_parts))
        return "\n\n".join(parts)

    def _normalize_text(self, text: str) -> str:
        lines = [line.rstrip() for line in text.replace("\r\n", "\n").replace("\r", "\n").split("\n")]
        normalized: list[str] = []
        blank = False
        for line in lines:
            stripped = line.strip()
            if not stripped:
                if not blank:
                    normalized.append("")
                blank = True
                continue
            normalized.append(stripped)
            blank = False
        return "\n".join(normalized).strip()

    def _map_name_for_param(self) -> str:
        if self.extracted_data:
            return self.extracted_data.get("map", {}).get("mapname", "")
        return self._current_map_name

    def load_extracted_data(self, json_path: str) -> None:
        self.extracted_data = json.loads(Path(json_path).read_text(encoding="utf-8"))

    def build_skill(self):
        if not self.extracted_data:
            raise RuntimeError("No IKM extracted data available")

        skill_dir = Path(self.skill_dir)
        refs_dir = skill_dir / "references"
        raw_dir = skill_dir / "raw"
        refs_dir.mkdir(parents=True, exist_ok=True)
        raw_dir.mkdir(parents=True, exist_ok=True)

        (raw_dir / "map.json").write_text(
            json.dumps(self.extracted_data.get("map", {}), ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
        (raw_dir / "source.json").write_text(
            json.dumps(
                {
                    "mode": self.extracted_data.get("mode"),
                    "portal": self.extracted_data.get("portal"),
                    "search": self.extracted_data.get("search"),
                    "keyword": self.extracted_data.get("keyword"),
                },
                ensure_ascii=False,
                indent=2,
            ),
            encoding="utf-8",
        )
        (raw_dir / "assets.json").write_text(
            json.dumps(self.extracted_data.get("assets", []), ensure_ascii=False, indent=2),
            encoding="utf-8",
        )

        self._write_index(refs_dir / "index.md")
        self._write_assets(refs_dir / "assets.md")
        self._write_attachments(refs_dir / "attachments.md")
        self._write_content_index(refs_dir / "content.md")
        self._write_skill(skill_dir / "SKILL.md")
        return True

    def _write_index(self, path: Path) -> None:
        data = self.extracted_data or {}
        map_info = data.get("map", {})
        source_title = self._source_title()
        content = [
            f"# {source_title} IKM Knowledge",
            "",
            f"- Source: {self.base_url}",
            f"- Map PK: {map_info.get('pkMapview') or self.pk}",
            f"- Mode: {data.get('mode', self.mode)}",
            f"- Keyword: {data.get('keyword') or ''}",
            f"- Description: {map_info.get('mapdesc') or ''}",
            f"- Owner org: {map_info.get('operatororgname') or ''}",
            f"- Publish date: {map_info.get('publishdate') or ''}",
            f"- Assets: {data.get('total_assets', 0)}",
            "",
            "## Reference Files",
            "",
            "- `assets.md` - IKM asset metadata",
            "- `attachments.md` - Attachment metadata and download paths",
            "- `content.md` - Parsed attachment content index, when available",
            "",
        ]
        path.write_text("\n".join(content), encoding="utf-8")

    def _write_assets(self, path: Path) -> None:
        lines = ["# IKM Assets", ""]
        for asset in (self.extracted_data or {}).get("assets", []):
            lines.extend(
                [
                    f"## {asset.get('assetname') or asset.get('pkKnowledgeAsset')}",
                    "",
                    f"- Asset code: {asset.get('assetcode') or ''}",
                    f"- Asset PK: {asset.get('pkKnowledgeAsset') or ''}",
                    f"- Type: {asset.get('ownedtypename') or ''}",
                    f"- Rule: {asset.get('ownedrulename') or ''}",
                    f"- Owner: {asset.get('ownedername') or ''}",
                    f"- Views: {asset.get('views') or 0}",
                    f"- Downloads: {asset.get('downloads') or 0}",
                    f"- File count: {asset.get('filenum') or 0}",
                    "",
                    asset.get("knowledgeitemdes") or "",
                    "",
                    f"Attachment names: {asset.get('detailnames') or ''}",
                    "",
                ]
            )
        path.write_text("\n".join(lines), encoding="utf-8")

    def _write_attachments(self, path: Path) -> None:
        lines = ["# IKM Attachments", ""]
        for asset in (self.extracted_data or {}).get("assets", []):
            lines.extend([f"## {asset.get('assetname') or asset.get('pkKnowledgeAsset')}", ""])
            for detail in asset.get("attachments", []):
                lines.extend(
                    [
                        f"- Name: {_strip_html(detail.get('oridetailname') or detail.get('detailname'))}",
                        f"  PK: {detail.get('pkKnowledgeDetail') or ''}",
                        f"  Type: {detail.get('detailtype') or detail.get('detailTypeCode') or ''}",
                        f"  Apply state: {detail.get('applyState')}",
                        f"  Download path: {detail.get('download_path') or ''}",
                        f"  Download skipped: {detail.get('download_skipped') or ''}",
                        f"  Download error: {detail.get('download_error') or ''}",
                        f"  Content path: {detail.get('content_path') or ''}",
                        f"  Parse error: {detail.get('parse_error') or ''}",
                    ]
                )
            lines.append("")
        path.write_text("\n".join(lines), encoding="utf-8")

    def _write_content_index(self, path: Path) -> None:
        lines = ["# IKM Parsed Attachment Content", ""]
        has_content = False
        for asset in (self.extracted_data or {}).get("assets", []):
            asset_lines: list[str] = []
            for detail in asset.get("attachments", []):
                content_path = detail.get("content_path")
                if not content_path:
                    continue
                has_content = True
                title = detail.get("oridetailname") or detail.get("detailname") or content_path
                asset_lines.extend(
                    [
                        f"### {_strip_html(title)}",
                        "",
                        f"- Content file: `{content_path}`",
                        f"- Attachment PK: `{detail.get('pkKnowledgeDetail') or ''}`",
                        f"- Characters: {detail.get('content_chars') or 0}",
                        "",
                        detail.get("content_excerpt") or "",
                        "",
                    ]
                )
            if asset_lines:
                lines.extend([f"## {asset.get('assetname') or asset.get('pkKnowledgeAsset')}", ""])
                lines.extend(asset_lines)

        if has_content:
            path.write_text("\n".join(lines), encoding="utf-8")
        elif path.exists():
            path.unlink()

    def _write_skill(self, path: Path) -> None:
        data = self.extracted_data or {}
        map_info = data.get("map", {})
        source_title = self._source_title()
        has_parsed_content = any(
            detail.get("content_path")
            for asset in data.get("assets", [])
            for detail in asset.get("attachments", [])
        )
        description = self.config.get("description") or (
            f"Use when answering questions about the IKM {data.get('mode', self.mode)} export "
            f"{source_title}, its assets, "
            "parsed attachments, owners, categories, and reusable enterprise knowledge."
        )
        content_rule = (
            "Use `references/content.md` and `references/content/` first for substantive "
            "knowledge from parsed attachments."
            if has_parsed_content
            else "Attachment body text was not parsed; use metadata and downloaded files only."
        )
        content_reference = (
            "- `references/content.md` - parsed attachment content index\n"
            "- `references/content/` - parsed attachment body files"
            if has_parsed_content
            else "- `downloads/` - downloaded attachment files, when available"
        )
        content = f"""---
name: {self.name}
description: {description}
---

# {source_title} IKM Knowledge Skill

Use this skill when working with the local IKM {data.get('mode', self.mode)} export for **{source_title}**.

## Source

- IKM base URL: {self.base_url}
- Map PK: {map_info.get('pkMapview') or self.pk}
- Mode: {data.get('mode', self.mode)}
- Keyword: {data.get('keyword') or ''}
- Owner org: {map_info.get('operatororgname') or ''}
- Publish date: {map_info.get('publishdate') or ''}
- Assets extracted: {data.get('total_assets', 0)}

## How to Answer

- Use `references/index.md` first to understand the map.
- {content_rule}
- Use `references/assets.md` for asset metadata, descriptions, owners, and categories.
- Use `references/attachments.md` for attachment names, types, ids, and local download paths.
- Do not claim live IKM status unless the user explicitly asks to refresh the source.

## Reference Map

- `references/index.md` - overview
- `references/assets.md` - asset records
- `references/attachments.md` - attachment records
{content_reference}
- `raw/map.json` - raw map metadata
- `raw/source.json` - raw source mode metadata
- `raw/assets.json` - raw asset and attachment metadata
"""
        path.write_text(content, encoding="utf-8")

    def _source_title(self) -> str:
        data = self.extracted_data or {}
        map_info = data.get("map", {})
        if map_info.get("mapname"):
            return map_info["mapname"]
        if data.get("keyword"):
            return f"Search: {data['keyword']}"
        if data.get("portal", {}).get("actionlocid"):
            return f"Portal: {data['portal']['actionlocid']}"
        if data.get("asset", {}).get("pkKnowledgeAsset"):
            asset = data["asset"]
            return _strip_html(asset.get("assetname")) or f"Asset: {asset['pkKnowledgeAsset']}"
        return self.name


def main() -> int:
    parser = argparse.ArgumentParser(
        description="Extract iKM knowledge maps, assets, and attachments into a skill"
    )
    add_ikm_arguments(parser)
    args = parser.parse_args()
    converter = IKMToSkillConverter(vars(args))
    return converter.run()


if __name__ == "__main__":
    raise SystemExit(main())
