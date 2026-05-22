#!/usr/bin/env python3
"""
PowerPoint (.pptx) Presentation to Skill Converter

Converts PowerPoint presentations into AI-ready skills.
Uses python-pptx to extract slide content including text, tables, speaker notes,
images, and code blocks. Supports single files and directories of .pptx files.

Slides are grouped into sections based on layout type (section/title layouts act
as section breaks). Each section becomes a reference file in the output skill.

Usage:
    yonyou-doc2skill pptx --pptx presentation.pptx --name myskill
    yonyou-doc2skill pptx --pptx ./slides_dir/ --name myskill
    yonyou-doc2skill pptx --from-json presentation_extracted.json
"""

import json
import logging
import os
import re
from pathlib import Path

# Optional dependency guard
try:
    from pptx import Presentation
    from pptx.enum.text import PP_ALIGN  # noqa: F401
    from pptx.util import Emu  # noqa: F401

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from yonyou_doc2skill.cli.skill_converter import SkillConverter

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Monospace / code font families used for code-block detection
# ---------------------------------------------------------------------------
MONOSPACE_FONTS = frozenset(
    {
        "courier",
        "courier new",
        "consolas",
        "menlo",
        "monaco",
        "lucida console",
        "lucida sans typewriter",
        "dejavu sans mono",
        "liberation mono",
        "source code pro",
        "fira code",
        "fira mono",
        "jetbrains mono",
        "roboto mono",
        "ubuntu mono",
        "inconsolata",
        "hack",
        "cascadia code",
        "cascadia mono",
        "sf mono",
        "andale mono",
        "ibm plex mono",
        "droid sans mono",
        "noto mono",
        "pt mono",
        "overpass mono",
    }
)

# Layout names that typically signal a section/title divider slide
SECTION_LAYOUT_NAMES = frozenset(
    {
        "section header",
        "section",
        "title slide",
        "title only",
        "title and content",
        "blank",
    }
)

# Layout names that are strong section-break indicators (title-only slides)
TITLE_ONLY_LAYOUTS = frozenset(
    {
        "section header",
        "section",
        "title slide",
        "title only",
    }
)


def _check_pptx_deps() -> None:
    """Raise RuntimeError if python-pptx is not installed."""
    if not PPTX_AVAILABLE:
        raise RuntimeError(
            "python-pptx is required for PowerPoint support.\n"
            'Install with: pip install "yonyou-doc2skill[pptx]"\n'
            "Or: pip install python-pptx"
        )


def infer_description_from_pptx(
    metadata: dict | None = None,
    name: str = "",
) -> str:
    """Infer skill description from PowerPoint metadata or name.

    Tries to extract a meaningful description from:
    1. Presentation subject field
    2. Presentation title field
    3. Falls back to a template using the skill name

    Args:
        metadata: Presentation metadata dict with title, subject, author, etc.
        name: Skill name for fallback

    Returns:
        Description string suitable for "Use when..." format
    """
    if metadata:
        # Try subject field first (often contains a description)
        if metadata.get("subject"):
            desc = str(metadata["subject"]).strip()
            if len(desc) > 20:
                if len(desc) > 150:
                    desc = desc[:147] + "..."
                return f"Use when {desc.lower()}"

        # Try title if meaningful
        if metadata.get("title"):
            title = str(metadata["title"]).strip()
            if len(title) > 10 and not title.lower().endswith(".pptx"):
                return f"Use when working with {title.lower()}"

    return (
        f"Use when referencing {name} presentation"
        if name
        else "Use when referencing this presentation"
    )


# ---------------------------------------------------------------------------
# Main converter class
# ---------------------------------------------------------------------------


class PptxToSkillConverter(SkillConverter):
    """Convert PowerPoint presentation (.pptx) to an AI-ready skill.

    Follows the same pipeline pattern as the Word, EPUB, and PDF scrapers:
    extract -> categorize -> build_skill (reference files + index + SKILL.md).

    The extraction phase uses python-pptx to read slides, extracting:
    - Slide titles, body text, and speaker notes
    - Tables (converted to markdown)
    - Image counts and descriptions
    - Code blocks (detected via monospace font usage)
    - Presentation-level metadata (title, author, subject, etc.)
    - Slide layout information for section grouping

    Supports both single .pptx files and directories containing multiple
    .pptx files (merged into a single skill).
    """

    SOURCE_TYPE = "pptx"

    def __init__(self, config: dict) -> None:
        """Initialize the converter with a configuration dictionary.

        Args:
            config: Configuration dict with keys:
                - name (str): Skill name (required)
                - pptx_path (str): Path to .pptx file or directory (optional)
                - description (str): Skill description (optional, inferred if absent)
                - categories (dict): Manual category assignments (optional)
        """
        super().__init__(config)
        self.config = config
        self.name: str = config["name"]
        self.pptx_path: str = config.get("pptx_path", "")
        self.description: str = (
            config.get("description") or f"Use when referencing {self.name} presentation"
        )

        # Paths
        self.skill_dir: str = f"output/{self.name}"
        self.data_file: str = f"output/{self.name}_extracted.json"

        # Categories config
        self.categories: dict = config.get("categories", {})

        # Extracted data (populated by extract_pptx or load_extracted_data)
        self.extracted_data: dict | None = None

    def extract(self):
        """Extract content from PowerPoint files (SkillConverter interface)."""
        self.extract_pptx()

    # ------------------------------------------------------------------
    # Extraction
    # ------------------------------------------------------------------

    def extract_pptx(self) -> bool:
        """Extract content from PowerPoint file(s) using python-pptx.

        Handles both single .pptx files and directories containing multiple
        .pptx files. For directories, files are processed in sorted order and
        their slides are concatenated sequentially.

        Workflow:
        1. Check dependencies (python-pptx)
        2. Resolve input path (single file vs. directory)
        3. For each .pptx file:
           a. Open with python-pptx Presentation class
           b. Extract presentation-level metadata
           c. Iterate slides, extracting text, notes, tables, images, code
        4. Detect section breaks from slide layouts
        5. Group slides into sections
        6. Detect code languages via LanguageDetector
        7. Save intermediate JSON to {name}_extracted.json

        Returns:
            True on successful extraction.

        Raises:
            FileNotFoundError: If the pptx_path does not exist.
            ValueError: If no .pptx files are found in a directory.
            RuntimeError: If extraction fails for other reasons.
        """
        _check_pptx_deps()

        print(f"\n🔍 Extracting from PowerPoint: {self.pptx_path}")

        pptx_path = Path(self.pptx_path)
        if not pptx_path.exists():
            raise FileNotFoundError(f"PowerPoint path not found: {self.pptx_path}")

        # Collect .pptx file(s) to process
        pptx_files: list[Path] = []
        if pptx_path.is_dir():
            pptx_files = sorted(pptx_path.glob("*.pptx"))
            if not pptx_files:
                raise ValueError(f"No .pptx files found in directory: {self.pptx_path}")
            print(f"   Found {len(pptx_files)} .pptx file(s) in directory")
        else:
            if not str(pptx_path).lower().endswith(".pptx"):
                raise ValueError(f"Not a PowerPoint file (expected .pptx): {self.pptx_path}")
            pptx_files = [pptx_path]

        # Accumulate slides across all files
        all_slides: list[dict] = []
        merged_metadata: dict = {}
        total_image_count = 0
        slide_offset = 0

        for file_path in pptx_files:
            print(f"   Processing: {file_path.name}")
            try:
                prs = Presentation(str(file_path))
            except Exception as e:
                raise RuntimeError(f"Failed to open PowerPoint file: {file_path}\n{e}") from e

            # Extract metadata from first (or only) file
            if not merged_metadata:
                merged_metadata = self._extract_presentation_metadata(prs)
                if merged_metadata.get("title"):
                    print(f"   Title: {merged_metadata['title']}")
                if merged_metadata.get("author"):
                    print(f"   Author: {merged_metadata['author']}")

            # Extract each slide
            for slide_idx, slide in enumerate(prs.slides):
                slide_number = slide_offset + slide_idx + 1
                slide_data = self._extract_slide(slide, slide_number)

                # Track source file for multi-file scenarios
                if len(pptx_files) > 1:
                    slide_data["source_file"] = file_path.name

                all_slides.append(slide_data)
                total_image_count += slide_data.get("image_count", 0)

            slide_offset += len(prs.slides)

        print(f"   Total slides extracted: {len(all_slides)}")

        # Update description from metadata if not explicitly set
        if not self.config.get("description"):
            self.description = infer_description_from_pptx(merged_metadata, self.name)

        # Group slides into sections based on layout and section breaks
        sections = self._group_slides_into_sections(all_slides)

        # Detect code languages using LanguageDetector
        languages_detected, total_code_blocks = self._detect_languages(sections)

        # Count total tables
        total_tables = sum(
            len(slide.get("tables", []))
            for section in sections
            for slide in section.get("slides", [])
        )

        result_data = {
            "source_file": self.pptx_path,
            "metadata": merged_metadata,
            "total_slides": len(all_slides),
            "total_sections": len(sections),
            "total_code_blocks": total_code_blocks,
            "total_images": total_image_count,
            "total_tables": total_tables,
            "languages_detected": languages_detected,
            "pages": sections,  # "pages" key for pipeline compatibility
        }

        # Save extracted data
        os.makedirs(os.path.dirname(self.data_file) or ".", exist_ok=True)
        with open(self.data_file, "w", encoding="utf-8") as f:
            json.dump(result_data, f, indent=2, ensure_ascii=False, default=str)

        print(f"\n💾 Saved extracted data to: {self.data_file}")
        self.extracted_data = result_data
        print(
            f"✅ Extracted {len(sections)} sections ({len(all_slides)} slides), "
            f"{total_code_blocks} code blocks, "
            f"{total_image_count} images, "
            f"{total_tables} tables"
        )
        return True

    def _extract_presentation_metadata(self, prs) -> dict:
        """Extract presentation-level metadata from core properties.

        Reads the Office Open XML core properties: title, author, subject,
        category, comments, keywords, created/modified dates, revision,
        and last_modified_by.

        Args:
            prs: python-pptx Presentation object

        Returns:
            Dictionary of metadata fields (string values, None for missing).
        """
        props = prs.core_properties
        return {
            "title": props.title or "",
            "author": props.author or "",
            "subject": props.subject or "",
            "category": props.category or "",
            "comments": props.comments or "",
            "keywords": props.keywords or "",
            "created": str(props.created) if props.created else "",
            "modified": str(props.modified) if props.modified else "",
            "last_modified_by": props.last_modified_by or "",
            "revision": props.revision if props.revision else None,
            "slide_count": len(prs.slides),
            "slide_width": prs.slide_width,
            "slide_height": prs.slide_height,
        }

    def _extract_slide(self, slide, slide_number: int) -> dict:
        """Extract all content from a single slide.

        Processes the slide's shapes to extract:
        - Title text (from the title placeholder)
        - Body text (from all text frames, excluding title)
        - Speaker notes
        - Tables (as structured data)
        - Image count and descriptions
        - Code blocks (detected by monospace font usage)
        - Layout name and type information

        Args:
            slide: python-pptx Slide object
            slide_number: 1-based slide number in the presentation

        Returns:
            Dictionary with all extracted slide data.
        """
        layout_name = ""
        if slide.slide_layout:
            layout_name = slide.slide_layout.name or ""

        # Determine if this is a section/title slide
        is_section_slide = layout_name.lower() in TITLE_ONLY_LAYOUTS

        # Extract title
        title = ""
        if slide.shapes.title:
            title = slide.shapes.title.text.strip()

        # Extract body text from all text frames (excluding title placeholder)
        body_parts: list[str] = []
        code_blocks: list[dict] = []
        image_count = 0
        tables: list[dict] = []

        for shape in slide.shapes:
            # Skip the title placeholder (already extracted)
            if shape.has_text_frame and shape == slide.shapes.title:
                continue

            # Process grouped shapes recursively
            if shape.shape_type is not None and hasattr(shape, "shapes"):
                group_text, group_codes, group_images = self._extract_group_shapes(shape)
                body_parts.extend(group_text)
                code_blocks.extend(group_codes)
                image_count += group_images
                continue

            # Tables
            if shape.has_table:
                table_data = self._extract_tables(shape.table)
                if table_data:
                    tables.append(table_data)
                continue

            # Images
            if self._is_image_shape(shape):
                image_count += 1
                continue

            # Text frames
            if shape.has_text_frame:
                frame_text, frame_codes = self._process_text_frame(shape.text_frame)
                if frame_codes:
                    code_blocks.extend(frame_codes)
                elif frame_text:
                    body_parts.append(frame_text)

        # Extract speaker notes
        speaker_notes = self._extract_speaker_notes(slide)

        # Extract image info summary
        image_info = self._extract_images_info(slide)

        return {
            "slide_number": slide_number,
            "layout_name": layout_name,
            "is_section_slide": is_section_slide,
            "title": title,
            "body_text": "\n\n".join(body_parts),
            "speaker_notes": speaker_notes,
            "tables": tables,
            "code_blocks": code_blocks,
            "image_count": image_count,
            "image_info": image_info,
        }

    def _extract_group_shapes(self, group_shape) -> tuple[list[str], list[dict], int]:
        """Recursively extract content from grouped shapes.

        PowerPoint allows shapes to be grouped together. This method walks
        the group hierarchy and extracts text, code blocks, and image counts
        from all nested shapes.

        Args:
            group_shape: python-pptx GroupShape object

        Returns:
            Tuple of (text_parts, code_blocks, image_count)
        """
        text_parts: list[str] = []
        code_blocks: list[dict] = []
        image_count = 0

        for shape in group_shape.shapes:
            # Nested groups
            if hasattr(shape, "shapes"):
                sub_text, sub_codes, sub_images = self._extract_group_shapes(shape)
                text_parts.extend(sub_text)
                code_blocks.extend(sub_codes)
                image_count += sub_images
                continue

            # Tables in groups
            if shape.has_table:
                # Tables in groups are rare but possible; skip for text extraction
                continue

            # Images in groups
            if self._is_image_shape(shape):
                image_count += 1
                continue

            # Text frames in groups
            if shape.has_text_frame:
                frame_text, frame_codes = self._process_text_frame(shape.text_frame)
                if frame_codes:
                    code_blocks.extend(frame_codes)
                elif frame_text:
                    text_parts.append(frame_text)

        return text_parts, code_blocks, image_count

    def _process_text_frame(self, text_frame) -> tuple[str, list[dict]]:
        """Process a text frame, separating regular text from code blocks.

        Examines the font properties of each paragraph's runs to determine
        whether the content is code (monospace font) or regular text.

        Args:
            text_frame: python-pptx TextFrame object

        Returns:
            Tuple of (plain_text, code_blocks) where code_blocks is a list
            of dicts with 'code', 'language', and 'quality_score' keys.
        """
        text_parts: list[str] = []
        code_parts: list[str] = []
        code_blocks: list[dict] = []
        in_code_block = False

        for paragraph in text_frame.paragraphs:
            para_text = paragraph.text.strip()
            if not para_text:
                # Empty paragraph may separate code blocks
                if in_code_block and code_parts:
                    code_blocks.append(self._finalize_code_block(code_parts))
                    code_parts = []
                    in_code_block = False
                continue

            is_code = self._detect_code_blocks(paragraph)

            if is_code:
                in_code_block = True
                code_parts.append(paragraph.text)
            else:
                # Flush any accumulated code
                if in_code_block and code_parts:
                    code_blocks.append(self._finalize_code_block(code_parts))
                    code_parts = []
                    in_code_block = False
                text_parts.append(para_text)

        # Flush trailing code block
        if code_parts:
            code_blocks.append(self._finalize_code_block(code_parts))

        return "\n".join(text_parts), code_blocks

    def _finalize_code_block(self, code_parts: list[str]) -> dict:
        """Create a code block dict from accumulated code lines.

        Args:
            code_parts: List of code line strings

        Returns:
            Dict with 'code', 'language', and 'quality_score' keys.
        """
        code_text = "\n".join(code_parts)
        quality = _score_code_quality(code_text)
        return {
            "code": code_text,
            "language": "",
            "quality_score": quality,
        }

    def _extract_tables(self, table) -> dict | None:
        """Extract table data from a python-pptx Table object.

        Converts the table into a structured dict with headers and rows.
        The first row is treated as the header row.

        Args:
            table: python-pptx Table object

        Returns:
            Dict with 'headers' (list[str]) and 'rows' (list[list[str]]) keys,
            or None if the table is empty.
        """
        if not table.rows:
            return None

        rows_data: list[list[str]] = []
        for row in table.rows:
            cells = []
            for cell in row.cells:
                # Extract text from all paragraphs in the cell
                cell_text = "\n".join(p.text.strip() for p in cell.text_frame.paragraphs).strip()
                cells.append(cell_text)
            rows_data.append(cells)

        if not rows_data:
            return None

        # First row is headers
        headers = rows_data[0]
        data_rows = rows_data[1:]

        return {"headers": headers, "rows": data_rows}

    def _extract_images_info(self, slide) -> list[dict]:
        """Extract descriptive information about images on a slide.

        Does not extract image binary data (to keep JSON output manageable).
        Instead, records image position, size, and any alt text or name.

        Args:
            slide: python-pptx Slide object

        Returns:
            List of dicts with image metadata (name, width, height, alt_text).
        """
        images: list[dict] = []

        for shape in slide.shapes:
            if not self._is_image_shape(shape):
                continue

            info: dict = {
                "index": len(images),
                "name": shape.name or "",
                "width": shape.width if hasattr(shape, "width") else 0,
                "height": shape.height if hasattr(shape, "height") else 0,
            }

            # Try to get alt text (accessibility description)
            try:
                # python-pptx stores alt text in shape._element
                desc_elem = shape._element.find(
                    ".//{http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing}cNvPr"
                )
                if desc_elem is not None:
                    info["alt_text"] = desc_elem.get("descr", "")
                else:
                    # Try the main namespace
                    for child in shape._element.iter():
                        descr = child.get("descr")
                        if descr:
                            info["alt_text"] = descr
                            break
            except Exception:
                pass

            images.append(info)

        return images

    def _detect_code_blocks(self, paragraph) -> bool:
        """Detect whether a paragraph contains code based on font properties.

        Code blocks in presentations are typically identified by:
        1. Monospace font family (Courier, Consolas, etc.)
        2. Small font size relative to body text
        3. Specific formatting patterns (e.g., syntax-highlighted runs)

        This method checks the font properties of the paragraph's runs
        and uses heuristics to determine if the content is code.

        Args:
            paragraph: python-pptx Paragraph object

        Returns:
            True if the paragraph appears to contain code.
        """
        if not paragraph.runs:
            return False

        # Count runs with monospace fonts
        mono_runs = 0
        total_runs = 0
        total_chars = 0
        mono_chars = 0

        for run in paragraph.runs:
            run_text = run.text
            if not run_text.strip():
                continue

            total_runs += 1
            char_count = len(run_text)
            total_chars += char_count

            font_name = ""
            if run.font and run.font.name:
                font_name = run.font.name.lower()

            if font_name in MONOSPACE_FONTS:
                mono_runs += 1
                mono_chars += char_count

        if total_runs == 0 or total_chars == 0:
            return False

        # If majority of characters are in monospace font, it's code
        mono_ratio = mono_chars / total_chars
        if mono_ratio >= 0.6:
            return True

        # Also check the paragraph text for code-like patterns
        text = paragraph.text.strip()
        return mono_ratio >= 0.3 and self._text_looks_like_code(text)

    def _text_looks_like_code(self, text: str) -> bool:
        """Heuristic check whether text content looks like source code.

        Uses pattern matching to detect common code constructs like
        function definitions, imports, variable assignments, etc.

        Args:
            text: The text content to check

        Returns:
            True if the text appears to be source code.
        """
        if not text:
            return False

        # Strong code indicators
        code_patterns = [
            r"^\s*(def |class |function |func |fn |pub fn )",
            r"^\s*(import |from .+ import|require\(|#include|using )",
            r"^\s*(if |else:|elif |for |while |switch |case )",
            r"^\s*(return |yield |raise |throw )",
            r"^\s*(const |let |var |int |float |str |bool )",
            r"[{}\[\]();]",
            r"^\s*#\s*\w+",  # preprocessor or comment
            r"=>|->|\|\||&&",  # operators
            r"^\s*@\w+",  # decorators
            r'^\s*\w+\s*=\s*["\'\d\[\{]',  # assignment
            r"^\s*\$\w+",  # shell/PHP variables
            r"^\s*(SELECT|INSERT|UPDATE|DELETE|CREATE|FROM|WHERE)\s",  # SQL
        ]

        for pattern in code_patterns:
            if re.search(pattern, text, re.MULTILINE | re.IGNORECASE):
                return True

        # Check ratio of special characters (code tends to have more)
        if len(text) > 10:
            special_count = sum(1 for c in text if c in "{}[]();=<>|&!@#$%^*~`")
            if special_count / len(text) > 0.08:
                return True

        return False

    def _extract_speaker_notes(self, slide) -> str:
        """Extract speaker notes from a slide.

        Speaker notes are stored in the slide's notes_slide object.
        Returns the full text of the notes, or empty string if none exist.

        Args:
            slide: python-pptx Slide object

        Returns:
            Speaker notes text string.
        """
        try:
            if not slide.has_notes_slide:
                return ""

            notes_slide = slide.notes_slide
            if not notes_slide or not notes_slide.notes_text_frame:
                return ""

            notes_text = notes_slide.notes_text_frame.text.strip()
            return notes_text
        except Exception:
            logger.debug(f"Could not extract speaker notes from slide {slide.slide_id}")
            return ""

    def _is_image_shape(self, shape) -> bool:
        """Check if a shape is an image (picture).

        Args:
            shape: python-pptx Shape object

        Returns:
            True if the shape contains an image.
        """
        try:
            # python-pptx shape_type 13 = PICTURE
            if (
                hasattr(shape, "shape_type")
                and shape.shape_type is not None
                and shape.shape_type == 13  # MSO_SHAPE_TYPE.PICTURE
            ):
                return True
            # Also check for image in the shape's element
            if hasattr(shape, "image"):
                return True
        except Exception:
            pass
        return False

    # ------------------------------------------------------------------
    # Section grouping
    # ------------------------------------------------------------------

    def _group_slides_into_sections(self, slides: list[dict]) -> list[dict]:
        """Group slides into sections based on layout type and section breaks.

        Section breaks are detected from:
        1. Slides with section/title-only layouts (is_section_slide=True)
        2. Slides whose title matches common section patterns

        Each section contains:
        - section_number: 1-based index
        - heading: Section title (from the section break slide)
        - heading_level: 'h1' for sections, 'h2' for subsections
        - text: Combined body text from all slides in the section
        - slides: List of raw slide dicts
        - code_samples: Aggregated code blocks
        - tables: Aggregated tables
        - speaker_notes: Combined speaker notes
        - image_count: Total images in the section

        Args:
            slides: List of slide dicts from _extract_slide()

        Returns:
            List of section dicts compatible with the pipeline format.
        """
        if not slides:
            return []

        # Identify section break points
        section_breaks: list[int] = []
        for i, slide in enumerate(slides):
            if slide.get("is_section_slide") and slide.get("title"):
                section_breaks.append(i)

        # If no explicit section breaks, treat the entire presentation as one section
        if not section_breaks:
            section = self._build_section_from_slides(
                section_number=1,
                heading=slides[0].get("title", self.name),
                heading_level="h1",
                slide_list=slides,
            )
            return [section]

        # Build sections from break points
        sections: list[dict] = []
        section_number = 0

        # Handle slides before the first section break
        if section_breaks[0] > 0:
            pre_section_slides = slides[: section_breaks[0]]
            section_number += 1
            section = self._build_section_from_slides(
                section_number=section_number,
                heading=pre_section_slides[0].get("title", "Introduction"),
                heading_level="h1",
                slide_list=pre_section_slides,
            )
            sections.append(section)

        # Process each section
        for idx, break_idx in enumerate(section_breaks):
            section_number += 1
            section_slide = slides[break_idx]
            heading = section_slide.get("title", f"Section {section_number}")

            # Determine end of this section
            end_idx = section_breaks[idx + 1] if idx + 1 < len(section_breaks) else len(slides)

            section_slides = slides[break_idx:end_idx]

            section = self._build_section_from_slides(
                section_number=section_number,
                heading=heading,
                heading_level="h1",
                slide_list=section_slides,
            )
            sections.append(section)

        return sections

    def _build_section_from_slides(
        self,
        section_number: int,
        heading: str,
        heading_level: str,
        slide_list: list[dict],
    ) -> dict:
        """Aggregate multiple slides into a single section dict.

        Combines text, code blocks, tables, and notes from all slides
        in the section into a single section dict compatible with the
        pipeline's intermediate JSON format.

        Args:
            section_number: 1-based section index
            heading: Section heading text
            heading_level: 'h1' or 'h2'
            slide_list: List of slide dicts to include

        Returns:
            Section dict with aggregated content.
        """
        text_parts: list[str] = []
        code_samples: list[dict] = []
        all_tables: list[dict] = []
        notes_parts: list[str] = []
        image_count = 0
        sub_headings: list[dict] = []

        for slide in slide_list:
            slide_num = slide.get("slide_number", "?")
            slide_title = slide.get("title", "")

            # Add slide title as sub-heading (unless it's the section heading)
            if slide_title and slide_title != heading:
                sub_headings.append(
                    {
                        "level": "h3",
                        "text": f"Slide {slide_num}: {slide_title}",
                    }
                )

            # Collect body text
            body = slide.get("body_text", "").strip()
            if body:
                text_parts.append(body)

            # Collect code blocks
            code_blocks = slide.get("code_blocks", [])
            code_samples.extend(code_blocks)

            # Collect tables
            tables = slide.get("tables", [])
            all_tables.extend(tables)

            # Collect speaker notes
            notes = slide.get("speaker_notes", "").strip()
            if notes:
                notes_parts.append(f"[Slide {slide_num}] {notes}")

            # Count images
            image_count += slide.get("image_count", 0)

        # Combine text with speaker notes appended
        combined_text = "\n\n".join(text_parts)
        if notes_parts:
            combined_text += "\n\n### Speaker Notes\n\n" + "\n\n".join(notes_parts)

        return {
            "section_number": section_number,
            "heading": heading,
            "heading_level": heading_level,
            "text": combined_text,
            "headings": sub_headings,
            "code_samples": code_samples,
            "tables": all_tables,
            "slides": slide_list,
            "image_count": image_count,
            "slide_range": (
                f"{slide_list[0]['slide_number']}-{slide_list[-1]['slide_number']}"
                if slide_list
                else ""
            ),
        }

    # ------------------------------------------------------------------
    # Language detection
    # ------------------------------------------------------------------

    def _detect_languages(
        self,
        sections: list[dict],
    ) -> tuple[dict[str, int], int]:
        """Detect programming languages in code blocks across all sections.

        Uses the project's LanguageDetector for automatic language detection
        when the language is not already set.

        Args:
            sections: List of section dicts with code_samples

        Returns:
            Tuple of (languages_detected dict, total_code_blocks count)
        """
        try:
            from yonyou_doc2skill.cli.language_detector import LanguageDetector

            detector = LanguageDetector(min_confidence=0.15)
        except ImportError:
            detector = None
            logger.debug("LanguageDetector not available, skipping language detection")

        languages_detected: dict[str, int] = {}
        total_code_blocks = 0

        for section in sections:
            for code_sample in section.get("code_samples", []):
                total_code_blocks += 1
                lang = code_sample.get("language", "")

                if lang:
                    languages_detected[lang] = languages_detected.get(lang, 0) + 1
                elif detector:
                    code = code_sample.get("code", "")
                    if code:
                        detected_lang, confidence = detector.detect_from_code(code)
                        if detected_lang and confidence >= 0.3:
                            code_sample["language"] = detected_lang
                            languages_detected[detected_lang] = (
                                languages_detected.get(detected_lang, 0) + 1
                            )

        return languages_detected, total_code_blocks

    # ------------------------------------------------------------------
    # Load / Categorize / Build
    # ------------------------------------------------------------------

    def load_extracted_data(self, json_path: str) -> bool:
        """Load previously extracted data from JSON file.

        Args:
            json_path: Path to the extracted JSON file

        Returns:
            True on success.
        """
        print(f"\n📂 Loading extracted data from: {json_path}")
        with open(json_path, encoding="utf-8") as f:
            self.extracted_data = json.load(f)
        total = self.extracted_data.get("total_sections", len(self.extracted_data.get("pages", [])))
        print(f"✅ Loaded {total} sections")
        return True

    def categorize_content(self) -> dict[str, dict]:
        """Categorize sections based on headings, keywords, or config.

        For a single PowerPoint source, creates one category containing all
        sections. For keyword-based categorization (multi-source), scores
        each section against category keywords.

        Returns:
            Dict mapping category keys to category dicts with 'title' and
            'pages' (list of sections).
        """
        print("\n📋 Categorizing content...")

        categorized: dict[str, dict] = {}
        sections = self.extracted_data.get("pages", [])

        # For single PPTX source, use single category with all sections
        if self.pptx_path:
            pptx_basename = Path(self.pptx_path).stem
            category_key = self._sanitize_filename(pptx_basename)
            categorized[category_key] = {
                "title": pptx_basename,
                "pages": sections,
            }
            print("✅ Created 1 category (single PowerPoint source)")
            print(f"   - {pptx_basename}: {len(sections)} sections")
            return categorized

        # Keyword-based categorization (multi-source scenario)
        if self.categories:
            first_value = next(iter(self.categories.values()), None)
            if isinstance(first_value, list) and first_value and isinstance(first_value[0], dict):
                # Already categorized format
                for cat_key, pages in self.categories.items():
                    categorized[cat_key] = {
                        "title": cat_key.replace("_", " ").title(),
                        "pages": pages,
                    }
            else:
                # Keyword-based categorization
                for cat_key in self.categories:
                    categorized[cat_key] = {
                        "title": cat_key.replace("_", " ").title(),
                        "pages": [],
                    }

                for section in sections:
                    text = section.get("text", "").lower()
                    heading_text = section.get("heading", "").lower()

                    scores: dict[str, int] = {}
                    for cat_key, keywords in self.categories.items():
                        if isinstance(keywords, list):
                            score = sum(
                                1
                                for kw in keywords
                                if isinstance(kw, str)
                                and (kw.lower() in text or kw.lower() in heading_text)
                            )
                        else:
                            score = 0
                        if score > 0:
                            scores[cat_key] = score

                    if scores:
                        best_cat = max(scores, key=scores.get)
                        categorized[best_cat]["pages"].append(section)
                    else:
                        if "other" not in categorized:
                            categorized["other"] = {"title": "Other", "pages": []}
                        categorized["other"]["pages"].append(section)
        else:
            # No categorization - single category
            categorized["content"] = {"title": "Content", "pages": sections}

        print(f"✅ Created {len(categorized)} categories")
        for _cat_key, cat_data in categorized.items():
            print(f"   - {cat_data['title']}: {len(cat_data['pages'])} sections")

        return categorized

    def build_skill(self) -> None:
        """Build complete skill structure from extracted data.

        Creates the output directory structure with:
        - references/ — one markdown file per category
        - references/index.md — category index with statistics
        - SKILL.md — main skill file with frontmatter and overview
        - scripts/ — empty (reserved for future use)
        - assets/ — empty (reserved for image export)
        """
        print(f"\n🏗️  Building skill: {self.name}")

        # Create directories
        os.makedirs(f"{self.skill_dir}/references", exist_ok=True)
        os.makedirs(f"{self.skill_dir}/scripts", exist_ok=True)
        os.makedirs(f"{self.skill_dir}/assets", exist_ok=True)

        # Categorize content
        categorized = self.categorize_content()

        # Generate reference files
        print("\n📝 Generating reference files...")
        total_sections = len(categorized)
        section_num = 1
        for cat_key, cat_data in categorized.items():
            self._generate_reference_file(cat_key, cat_data, section_num, total_sections)
            section_num += 1

        # Generate index
        self._generate_index(categorized)

        # Generate SKILL.md
        self._generate_skill_md(categorized)

        print(f"\n✅ Skill built successfully: {self.skill_dir}/")
        print(f"\n📦 Next step: Package with: yonyou-doc2skill package {self.skill_dir}/")

    # ------------------------------------------------------------------
    # Output generation (private)
    # ------------------------------------------------------------------

    def _generate_reference_file(
        self,
        _cat_key: str,
        cat_data: dict,
        section_num: int,
        total_sections: int,
    ) -> None:
        """Generate a reference markdown file for a category of sections.

        Each section's slides are rendered as markdown with slide numbers,
        body text, code examples, tables, speaker notes, and image counts.

        Args:
            _cat_key: Category key (unused, for interface consistency)
            cat_data: Category dict with 'title' and 'pages' keys
            section_num: 1-based index among all categories
            total_sections: Total number of categories being generated
        """
        sections = cat_data["pages"]

        # Use pptx basename for filename
        pptx_basename = ""
        if self.pptx_path:
            pptx_basename = Path(self.pptx_path).stem

        if sections:
            section_nums = [s.get("section_number", i + 1) for i, s in enumerate(sections)]

            if total_sections == 1:
                filename = (
                    f"{self.skill_dir}/references/{pptx_basename}.md"
                    if pptx_basename
                    else f"{self.skill_dir}/references/main.md"
                )
            else:
                sec_range = f"s{min(section_nums)}-s{max(section_nums)}"
                base_name = pptx_basename if pptx_basename else "section"
                filename = f"{self.skill_dir}/references/{base_name}_{sec_range}.md"
        else:
            filename = f"{self.skill_dir}/references/section_{section_num:02d}.md"

        with open(filename, "w", encoding="utf-8") as f:
            f.write(f"# {cat_data['title']}\n\n")

            for section in sections:
                sec_num = section.get("section_number", "?")
                heading = section.get("heading", "")
                heading_level = section.get("heading_level", "h1")
                slide_range = section.get("slide_range", "")

                f.write(f"---\n\n**📄 Source: Section {sec_num}**")
                if slide_range:
                    f.write(f" (Slides {slide_range})")
                f.write("\n\n")

                # Section heading
                if heading:
                    md_level = "#" * (int(heading_level[1]) + 1) if heading_level else "##"
                    f.write(f"{md_level} {heading}\n\n")

                # Sub-headings (individual slide titles)
                for sub_heading in section.get("headings", []):
                    sub_level = sub_heading.get("level", "h3")
                    sub_text = sub_heading.get("text", "")
                    if sub_text:
                        sub_md = "#" * (int(sub_level[1]) + 1) if sub_level else "###"
                        f.write(f"{sub_md} {sub_text}\n\n")

                # Body text
                text = section.get("text", "").strip()
                if text:
                    f.write(f"{text}\n\n")

                # Code samples
                code_list = section.get("code_samples", [])
                if code_list:
                    f.write("### Code Examples\n\n")
                    for code in code_list:
                        lang = code.get("language", "")
                        f.write(f"```{lang}\n{code['code']}\n```\n\n")

                # Tables as markdown
                tables = section.get("tables", [])
                if tables:
                    f.write("### Tables\n\n")
                    for table in tables:
                        headers = table.get("headers", [])
                        rows = table.get("rows", [])
                        if headers:
                            f.write("| " + " | ".join(str(h) for h in headers) + " |\n")
                            f.write("| " + " | ".join("---" for _ in headers) + " |\n")
                        for row in rows:
                            f.write("| " + " | ".join(str(c) for c in row) + " |\n")
                        f.write("\n")

                # Image count summary
                img_count = section.get("image_count", 0)
                if img_count > 0:
                    f.write(f"### Images\n\n*{img_count} image(s) in this section*\n\n")

                f.write("---\n\n")

        print(f"   Generated: {filename}")

    def _generate_index(self, categorized: dict[str, dict]) -> None:
        """Generate reference index file listing all categories and statistics.

        Args:
            categorized: Dict of category key -> category data
        """
        filename = f"{self.skill_dir}/references/index.md"

        pptx_basename = ""
        if self.pptx_path:
            pptx_basename = Path(self.pptx_path).stem

        total_sections = len(categorized)

        with open(filename, "w", encoding="utf-8") as f:
            f.write(f"# {self.name.title()} Presentation Reference\n\n")
            f.write("## Categories\n\n")

            section_num = 1
            for _cat_key, cat_data in categorized.items():
                sections = cat_data["pages"]
                section_count = len(sections)

                if sections:
                    section_nums = [s.get("section_number", i + 1) for i, s in enumerate(sections)]
                    sec_range_str = f"Sections {min(section_nums)}-{max(section_nums)}"

                    if total_sections == 1:
                        link_filename = f"{pptx_basename}.md" if pptx_basename else "main.md"
                    else:
                        sec_range = f"s{min(section_nums)}-s{max(section_nums)}"
                        base_name = pptx_basename if pptx_basename else "section"
                        link_filename = f"{base_name}_{sec_range}.md"
                else:
                    link_filename = f"section_{section_num:02d}.md"
                    sec_range_str = "N/A"

                f.write(
                    f"- [{cat_data['title']}]({link_filename}) "
                    f"({section_count} sections, {sec_range_str})\n"
                )
                section_num += 1

            f.write("\n## Statistics\n\n")
            f.write(f"- Total slides: {self.extracted_data.get('total_slides', 0)}\n")
            f.write(f"- Total sections: {self.extracted_data.get('total_sections', 0)}\n")
            f.write(f"- Code blocks: {self.extracted_data.get('total_code_blocks', 0)}\n")
            f.write(f"- Images: {self.extracted_data.get('total_images', 0)}\n")
            f.write(f"- Tables: {self.extracted_data.get('total_tables', 0)}\n")

            # Metadata
            metadata = self.extracted_data.get("metadata", {})
            if metadata.get("author"):
                f.write(f"- Author: {metadata['author']}\n")
            if metadata.get("created"):
                f.write(f"- Created: {metadata['created']}\n")

        print(f"   Generated: {filename}")

    def _generate_skill_md(self, categorized: dict[str, dict]) -> None:
        """Generate main SKILL.md file with YAML frontmatter and overview.

        Creates a comprehensive skill file with:
        - YAML frontmatter (name, description)
        - Document information (from metadata)
        - "When to Use" section
        - Section overview with slide counts
        - Key concepts from headings
        - Quick reference patterns
        - Top code examples grouped by language
        - Table summary
        - Documentation statistics
        - Navigation links

        Args:
            categorized: Dict of category key -> category data
        """
        filename = f"{self.skill_dir}/SKILL.md"

        skill_name = self.name.lower().replace("_", "-").replace(" ", "-")[:64]
        desc = self.description[:1024] if len(self.description) > 1024 else self.description

        with open(filename, "w", encoding="utf-8") as f:
            # YAML frontmatter
            f.write("---\n")
            f.write(f"name: {skill_name}\n")
            f.write(f"description: {desc}\n")
            f.write("---\n\n")

            f.write(f"# {self.name.title()} Presentation Skill\n\n")
            f.write(f"{self.description}\n\n")

            # Document metadata
            metadata = self.extracted_data.get("metadata", {})
            if any(v for v in metadata.values() if v):
                f.write("## 📋 Presentation Information\n\n")
                if metadata.get("title"):
                    f.write(f"**Title:** {metadata['title']}\n\n")
                if metadata.get("author"):
                    f.write(f"**Author:** {metadata['author']}\n\n")
                if metadata.get("subject"):
                    f.write(f"**Subject:** {metadata['subject']}\n\n")
                if metadata.get("category"):
                    f.write(f"**Category:** {metadata['category']}\n\n")
                if metadata.get("created"):
                    f.write(f"**Created:** {metadata['created']}\n\n")
                if metadata.get("modified"):
                    f.write(f"**Modified:** {metadata['modified']}\n\n")
                if metadata.get("slide_count"):
                    f.write(f"**Slides:** {metadata['slide_count']}\n\n")

            # When to Use
            f.write("## 💡 When to Use This Skill\n\n")
            f.write("Use this skill when you need to:\n")
            f.write(f"- Understand {self.name} concepts and fundamentals\n")
            f.write("- Review presentation content and key points\n")
            f.write("- Find code examples and implementation patterns\n")
            f.write("- Access speaker notes and additional context\n")
            f.write("- Reference tables and data from the presentation\n\n")

            # Section Overview
            total_slides = self.extracted_data.get("total_slides", 0)
            total_sections = self.extracted_data.get("total_sections", 0)
            f.write("## 📖 Section Overview\n\n")
            f.write(f"**Total Slides:** {total_slides}\n\n")
            f.write(f"**Total Sections:** {total_sections}\n\n")
            f.write("**Content Breakdown:**\n\n")
            for _cat_key, cat_data in categorized.items():
                section_count = len(cat_data["pages"])
                f.write(f"- **{cat_data['title']}**: {section_count} sections\n")
            f.write("\n")

            # Key Concepts from headings
            f.write(self._format_key_concepts())

            # Quick Reference patterns
            f.write("## ⚡ Quick Reference\n\n")
            f.write(self._format_patterns_from_content())

            # Code examples (top 15, grouped by language)
            all_code: list[dict] = []
            for section in self.extracted_data.get("pages", []):
                all_code.extend(section.get("code_samples", []))

            all_code.sort(key=lambda x: x.get("quality_score", 0), reverse=True)
            top_code = all_code[:15]

            if top_code:
                f.write("## 📝 Code Examples\n\n")
                f.write("*High-quality examples extracted from presentation*\n\n")

                by_lang: dict[str, list] = {}
                for code in top_code:
                    lang = code.get("language", "unknown")
                    by_lang.setdefault(lang, []).append(code)

                for lang in sorted(by_lang.keys()):
                    examples = by_lang[lang]
                    f.write(f"### {lang.title()} Examples ({len(examples)})\n\n")
                    for i, code in enumerate(examples[:5], 1):
                        quality = code.get("quality_score", 0)
                        code_text = code.get("code", "")
                        f.write(f"**Example {i}** (Quality: {quality:.1f}/10):\n\n")
                        f.write(f"```{lang}\n")
                        if len(code_text) <= 500:
                            f.write(code_text)
                        else:
                            f.write(code_text[:500] + "\n...")
                        f.write("\n```\n\n")

            # Table Summary (first 5 tables)
            all_tables: list[tuple[str, dict]] = []
            for section in self.extracted_data.get("pages", []):
                for table in section.get("tables", []):
                    all_tables.append((section.get("heading", ""), table))

            if all_tables:
                f.write("## 📊 Table Summary\n\n")
                f.write(f"*{len(all_tables)} table(s) found in presentation*\n\n")
                for section_heading, table in all_tables[:5]:
                    if section_heading:
                        f.write(f"**From section: {section_heading}**\n\n")
                    headers = table.get("headers", [])
                    rows = table.get("rows", [])
                    if headers:
                        f.write("| " + " | ".join(str(h) for h in headers) + " |\n")
                        f.write("| " + " | ".join("---" for _ in headers) + " |\n")
                        for row in rows[:5]:
                            f.write("| " + " | ".join(str(c) for c in row) + " |\n")
                        f.write("\n")

            # Statistics
            f.write("## 📊 Presentation Statistics\n\n")
            f.write(f"- **Total Slides**: {total_slides}\n")
            f.write(f"- **Total Sections**: {total_sections}\n")
            f.write(f"- **Code Blocks**: {self.extracted_data.get('total_code_blocks', 0)}\n")
            f.write(f"- **Images/Diagrams**: {self.extracted_data.get('total_images', 0)}\n")
            f.write(f"- **Tables**: {self.extracted_data.get('total_tables', 0)}\n")

            langs = self.extracted_data.get("languages_detected", {})
            if langs:
                f.write(f"- **Programming Languages**: {len(langs)}\n\n")
                f.write("**Language Breakdown:**\n\n")
                for lang, count in sorted(langs.items(), key=lambda x: x[1], reverse=True):
                    f.write(f"- {lang}: {count} examples\n")
                f.write("\n")

            # Navigation
            f.write("## 🗺️ Navigation\n\n")
            f.write("**Reference Files:**\n\n")
            for _cat_key, cat_data in categorized.items():
                cat_file = self._sanitize_filename(cat_data["title"])
                f.write(f"- `references/{cat_file}.md` - {cat_data['title']}\n")
            f.write("\n")
            f.write("See `references/index.md` for complete presentation structure.\n\n")

            # Footer
            f.write("---\n\n")
            f.write("**Generated by Skill Seeker** | PowerPoint Presentation Scraper\n")

        with open(filename, encoding="utf-8") as f:
            line_count = len(f.read().split("\n"))
        print(f"   Generated: {filename} ({line_count} lines)")

    # ------------------------------------------------------------------
    # Content analysis helpers
    # ------------------------------------------------------------------

    def _format_key_concepts(self) -> str:
        """Extract key concepts from section and slide headings.

        Returns:
            Markdown string with key concepts section, or empty string
            if no headings are found.
        """
        all_headings: list[tuple[str, str]] = []

        for section in self.extracted_data.get("pages", []):
            # Main section heading
            heading = section.get("heading", "").strip()
            level = section.get("heading_level", "h1")
            if heading and len(heading) > 3:
                all_headings.append((level, heading))
            # Sub-headings (individual slide titles)
            for sub in section.get("headings", []):
                text = sub.get("text", "").strip()
                sub_level = sub.get("level", "h3")
                if text and len(text) > 3:
                    all_headings.append((sub_level, text))

        if not all_headings:
            return ""

        content = "## 🔑 Key Concepts\n\n"
        content += "*Main topics covered in this presentation*\n\n"

        h1_headings = [text for level, text in all_headings if level == "h1"]
        h2_headings = [text for level, text in all_headings if level == "h2"]
        h3_headings = [text for level, text in all_headings if level == "h3"]

        if h1_headings:
            content += "**Major Sections:**\n\n"
            for heading in h1_headings[:10]:
                content += f"- {heading}\n"
            content += "\n"

        if h2_headings:
            content += "**Subsections:**\n\n"
            for heading in h2_headings[:15]:
                content += f"- {heading}\n"
            content += "\n"

        if h3_headings and not h2_headings:
            content += "**Slide Topics:**\n\n"
            for heading in h3_headings[:15]:
                content += f"- {heading}\n"
            content += "\n"

        return content

    def _format_patterns_from_content(self) -> str:
        """Extract common documentation patterns from section headings.

        Searches for keywords like "introduction", "overview", "demo",
        "agenda", etc. that are common in presentations.

        Returns:
            Markdown string describing found patterns.
        """
        patterns: list[dict] = []
        pattern_keywords = [
            "introduction",
            "overview",
            "agenda",
            "objectives",
            "getting started",
            "demo",
            "demonstration",
            "examples",
            "architecture",
            "design",
            "implementation",
            "best practices",
            "summary",
            "conclusion",
            "q&a",
            "questions",
            "next steps",
            "resources",
            "references",
            "appendix",
        ]

        for section in self.extracted_data.get("pages", []):
            heading_text = section.get("heading", "").lower()
            sec_num = section.get("section_number", 0)

            for keyword in pattern_keywords:
                if keyword in heading_text:
                    patterns.append(
                        {
                            "type": keyword.title(),
                            "heading": section.get("heading", ""),
                            "section": sec_num,
                        }
                    )
                    break

        if not patterns:
            return "*See reference files for detailed content*\n\n"

        content = "*Common presentation patterns found:*\n\n"
        by_type: dict[str, list] = {}
        for pattern in patterns:
            ptype = pattern["type"]
            by_type.setdefault(ptype, []).append(pattern)

        for ptype in sorted(by_type.keys()):
            items = by_type[ptype]
            content += f"**{ptype}** ({len(items)} sections):\n"
            for item in items[:3]:
                content += f"- {item['heading']} (section {item['section']})\n"
            content += "\n"

        return content

    def _sanitize_filename(self, name: str) -> str:
        """Convert a string to a filesystem-safe filename.

        Removes special characters, replaces spaces and hyphens with
        underscores, and lowercases the result.

        Args:
            name: Input string to sanitize

        Returns:
            Safe filename string
        """
        safe = re.sub(r"[^\w\s-]", "", name.lower())
        safe = re.sub(r"[-\s]+", "_", safe)
        return safe


# ---------------------------------------------------------------------------
# Module-level helpers
# ---------------------------------------------------------------------------


def _score_code_quality(code: str) -> float:
    """Score code quality on a 0-10 scale using heuristics.

    Higher scores indicate more substantial, well-structured code.
    Factors include line count, presence of definitions, imports,
    indentation, and code syntax characters.

    Args:
        code: Source code text to score

    Returns:
        Float quality score between 0.0 and 10.0
    """
    if not code:
        return 0.0

    score = 5.0
    lines = code.strip().split("\n")
    line_count = len(lines)

    # More lines = more substantial
    if line_count >= 10:
        score += 2.0
    elif line_count >= 5:
        score += 1.0

    # Has function/class definitions
    if re.search(r"\b(def |class |function |func |fn )", code):
        score += 1.5

    # Has imports/require
    if re.search(r"\b(import |from .+ import|require\(|#include|using )", code):
        score += 0.5

    # Has indentation (common in Python, JS, etc.)
    if re.search(r"^    ", code, re.MULTILINE):
        score += 0.5

    # Has assignment, operators, or common code syntax
    if re.search(r"[=:{}()\[\]]", code):
        score += 0.3

    # Very short snippets get penalized
    if len(code) < 30:
        score -= 2.0

    return min(10.0, max(0.0, score))
